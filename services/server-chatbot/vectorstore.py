import os
import json
import httpx
import tempfile
from datetime import datetime, timezone
from typing import Any, Optional

try:
    import fitz
except ImportError:  # pragma: no cover - exercised only when the optional runtime dependency is absent.
    fitz = None

from langchain_community.vectorstores import Chroma
from langchain_ollama import OllamaEmbeddings
from langchain_community.document_loaders import PyPDFLoader, TextLoader, Docx2txtLoader
from langchain_core.documents import Document
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_google_genai import GoogleGenerativeAIEmbeddings
from pptx import Presentation
from logger import get_logger

logger = get_logger(__name__)

CHROMA_DIR = os.getenv("CHROMA_DIR", "/app/chroma_db")

OLLAMA_BASE_URL = os.getenv("OLLAMA_BASE_URL", "http://ollama:11434")
EMBED_MODEL = os.getenv("EMBED_MODEL", "nomic-embed-text")

PLACEHOLDER_SECRETS = {"", "your-key-here", "ci-placeholder", "qa-compose-validation-placeholder"}

def normalise_optional_secret(value: str | None) -> str | None:
    """Treat local/CI placeholders as missing secrets before provider clients see them."""
    candidate = (value or "").strip()
    return None if candidate in PLACEHOLDER_SECRETS else candidate

USE_GEMINI_EMBEDDINGS = os.getenv("USE_GEMINI_EMBEDDINGS", "false").lower() == "true"
GEMINI_API_KEY = normalise_optional_secret(os.getenv("GEMINI_API_KEY"))
GEMINI_EMBED_MODEL = os.getenv("GEMINI_EMBED_MODEL", "models/gemini-embedding-001")

NODE_BASE_URL = os.getenv("NODE_BASE_URL", "http://server:4000")
SEARCH_MIN_RELEVANCE = float(os.getenv("SEARCH_MIN_RELEVANCE", "0.35"))
PDF_HIGHLIGHT_MAX_RECTS = int(os.getenv("PDF_HIGHLIGHT_MAX_RECTS", "24"))
PDF_HIGHLIGHT_MAX_SOURCE_CHARS = int(os.getenv("PDF_HIGHLIGHT_MAX_SOURCE_CHARS", "600"))

SOURCE_TYPE_ALIASES = {
    "infilenite": "resource",
    "file": "resource",
    "files": "resource",
    "resource": "resource",
    "resources": "resource",
    "convolution": "convolution_message",
    "chat": "convolution_message",
    "message": "convolution_message",
    "messages": "convolution_message",
    "annotation": "annotation",
    "annotations": "annotation",
    "coordidate": ["coordidate_session", "coordidate_poll"],
    "calendar": ["coordidate_session", "coordidate_poll"],
    "meeting": "coordidate_session",
    "meetings": "coordidate_session",
    "session": "coordidate_session",
    "sessions": "coordidate_session",
    "poll": "coordidate_poll",
    "polls": "coordidate_poll",
}

def _payload_get(item: Any, key: str, default: Any = None) -> Any:
    """Read pydantic models and dict payloads through the same small helper."""
    if isinstance(item, dict):
        return item.get(key, default)
    return getattr(item, key, default)

def _compact_text(value: Any, limit: int = 500) -> str:
    """Trim metadata snippets so source refs stay useful without bloating Chroma."""
    return " ".join(str(value or "").split())[:limit]

def _clean_metadata_value(value: Any) -> str | int | float | bool:
    """Chroma metadata must stay flat, so complex values are serialized."""
    if value is None:
        return ""
    if isinstance(value, (str, int, float, bool)):
        return value
    return json.dumps(value, ensure_ascii=False, sort_keys=True)

def _clean_metadata(metadata: dict[str, Any]) -> dict[str, str | int | float | bool]:
    """Drop empty keys and coerce metadata values into Chroma-compatible scalars."""
    cleaned: dict[str, str | int | float | bool] = {}
    for key, value in (metadata or {}).items():
        normalized_key = str(key or "").strip()
        if not normalized_key:
            continue
        cleaned[normalized_key] = _clean_metadata_value(value)
    return cleaned

def _iso_to_epoch_ms(value: Any) -> int:
    """Convert ISO-like date/time metadata into milliseconds for Chroma range filters."""
    text = str(value or "").strip()
    if not text:
        return 0
    try:
        normalized = text.replace("Z", "+00:00")
        if len(normalized) == 10:
            normalized = f"{normalized}T00:00:00+00:00"
        parsed = datetime.fromisoformat(normalized)
        if parsed.tzinfo is None:
            parsed = parsed.replace(tzinfo=timezone.utc)
        return int(parsed.timestamp() * 1000)
    except ValueError:
        return 0

def _source_label(source_ref: dict[str, Any], fallback: str = "") -> str:
    """Choose the label shown in source pills and tool summaries."""
    return str(source_ref.get("label") or source_ref.get("title") or fallback or "Domain source").strip()[:180]

def _pdf_scaled_rect(
    x0: Any,
    y0: Any,
    x1: Any,
    y1: Any,
    page_width: float,
    page_height: float,
    page_number: int,
) -> dict[str, Any] | None:
    """Convert a PyMuPDF word/line box into the highlighter's scaled position shape.

    react-pdf-highlighter-extended stores coordinates relative to a page size
    instead of the current zoom level. PyMuPDF's default page coordinates are
    also top-left based, so storing them with the page width/height lets the
    frontend scale them without brittle text matching.
    """
    try:
        left = max(0.0, min(float(x0), page_width))
        top = max(0.0, min(float(y0), page_height))
        right = max(0.0, min(float(x1), page_width))
        bottom = max(0.0, min(float(y1), page_height))
    except (TypeError, ValueError):
        return None

    if right <= left or bottom <= top or page_width <= 0 or page_height <= 0:
        return None

    return {
        "x1": round(left, 3),
        "y1": round(top, 3),
        "x2": round(right, 3),
        "y2": round(bottom, 3),
        "width": round(page_width, 3),
        "height": round(page_height, 3),
        "pageNumber": page_number,
    }

def _union_pdf_rects(rects: list[dict[str, Any]]) -> dict[str, Any] | None:
    """Build the bounding rectangle required by the PDF highlighter."""
    if not rects:
        return None
    page_number = int(rects[0].get("pageNumber") or 1)
    page_width = float(rects[0].get("width") or 0)
    page_height = float(rects[0].get("height") or 0)
    return {
        "x1": round(min(float(rect["x1"]) for rect in rects), 3),
        "y1": round(min(float(rect["y1"]) for rect in rects), 3),
        "x2": round(max(float(rect["x2"]) for rect in rects), 3),
        "y2": round(max(float(rect["y2"]) for rect in rects), 3),
        "width": round(page_width, 3),
        "height": round(page_height, 3),
        "pageNumber": page_number,
    }

def _valid_highlight_position(position: Any) -> dict[str, Any] | None:
    """Validate persisted highlight geometry before exposing it to the app."""
    if not isinstance(position, dict):
        return None
    bounding_rect = position.get("boundingRect")
    rects = position.get("rects")
    if not isinstance(bounding_rect, dict) or not isinstance(rects, list):
        return None
    return {
        "boundingRect": bounding_rect,
        "rects": [rect for rect in rects if isinstance(rect, dict)],
    }

class VectorStore:
    def __init__(self):
        if USE_GEMINI_EMBEDDINGS and GEMINI_API_KEY:
            logger.info(f"Using Gemini Embeddings | model: {GEMINI_EMBED_MODEL}")
            self.embeddings = GoogleGenerativeAIEmbeddings(
                model = GEMINI_EMBED_MODEL,
                google_api_key = GEMINI_API_KEY,
            )
        else:
            if USE_GEMINI_EMBEDDINGS and not GEMINI_API_KEY:
                logger.warning("Gemini embeddings requested without a real GEMINI_API_KEY; falling back to Ollama embeddings.")
            logger.info(f"Using Ollama Embeddings | model: {EMBED_MODEL} | url: {OLLAMA_BASE_URL}")
            self.embeddings = OllamaEmbeddings(
                model = EMBED_MODEL,
                base_url = OLLAMA_BASE_URL,
            )
        
        # The document text splitter, to find appropriate chunks
        self.splitter = RecursiveCharacterTextSplitter( 
            chunk_size = 500,
            chunk_overlap = 50,
        )
        
        self.db = Chroma(
            persist_directory = CHROMA_DIR,
            embedding_function = self.embeddings,
        )

    def _load_pptx_file(self, file_path: str, file_name: str) -> list[Document]:
        """Extract readable text from slides and tables in a PowerPoint deck."""
        presentation = Presentation(file_path)
        documents = []

        for index, slide in enumerate(presentation.slides, start=1):
            parts = []
            for shape in slide.shapes:
                if getattr(shape, "has_text_frame", False):
                    text = "\n".join(
                        paragraph.text.strip()
                        for paragraph in shape.text_frame.paragraphs
                        if paragraph.text and paragraph.text.strip()
                    )
                    if text:
                        parts.append(text)
                if getattr(shape, "has_table", False):
                    for row in shape.table.rows:
                        cells = [cell.text.strip() for cell in row.cells if cell.text and cell.text.strip()]
                        if cells:
                            parts.append(" | ".join(cells))

            page_content = "\n\n".join(parts).strip()
            if page_content:
                documents.append(
                    Document(
                        page_content=page_content,
                        metadata={"file_name": file_name, "slide": index},
                    )
                )

        if not documents:
            raise ValueError("PowerPoint file did not contain readable text")

        return documents

    def _extract_pdf_page_text(self, page: Any, page_number: int) -> tuple[str, list[dict[str, Any]], float, float]:
        """Extract page text plus word spans and boxes from a PDF page.

        The span offsets let us connect a LangChain text chunk back to the
        exact words it came from, which is more stable than asking the frontend
        to search the rendered PDF text layer after the model responds.
        """
        page_rect = page.rect
        page_width = float(page_rect.width or 0)
        page_height = float(page_rect.height or 0)
        raw_words = page.get_text("words", sort=True) or []

        parts: list[str] = []
        word_records: list[dict[str, Any]] = []
        cursor = 0
        last_line_key: tuple[int, int] | None = None

        for raw_word in raw_words:
            if len(raw_word) < 8:
                continue
            x0, y0, x1, y1, value, block_no, line_no, word_no = raw_word[:8]
            text = str(value or "").strip()
            if not text:
                continue

            line_key = (int(block_no), int(line_no))
            separator = "" if not parts else ("\n" if line_key != last_line_key else " ")
            if separator:
                parts.append(separator)
                cursor += len(separator)

            start = cursor
            parts.append(text)
            cursor += len(text)
            end = cursor
            last_line_key = line_key

            word_records.append({
                "start": start,
                "end": end,
                "line_key": (int(block_no), int(line_no), int(word_no)),
                "line_group": line_key,
                "x0": x0,
                "y0": y0,
                "x1": x1,
                "y1": y1,
                "page_number": page_number,
            })

        return "".join(parts).strip(), word_records, page_width, page_height

    def _pdf_rects_for_chunk(
        self,
        words: list[dict[str, Any]],
        start: int,
        end: int,
        page_width: float,
        page_height: float,
        page_number: int,
    ) -> list[dict[str, Any]]:
        """Merge chunk words into line rectangles suitable for text highlighting."""
        selected_words = [
            word for word in words
            if int(word.get("end") or 0) > start and int(word.get("start") or 0) < end
        ]
        if not selected_words:
            return []

        line_rects: list[dict[str, Any]] = []
        current_line: tuple[int, int] | None = None
        current_box: list[Any] | None = None

        for word in selected_words:
            line_group = word.get("line_group")
            if line_group != current_line:
                if current_box:
                    rect = _pdf_scaled_rect(*current_box, page_width, page_height, page_number)
                    if rect:
                        line_rects.append(rect)
                current_line = line_group
                current_box = [word.get("x0"), word.get("y0"), word.get("x1"), word.get("y1")]
                continue

            if current_box:
                current_box = [
                    min(float(current_box[0]), float(word.get("x0"))),
                    min(float(current_box[1]), float(word.get("y0"))),
                    max(float(current_box[2]), float(word.get("x1"))),
                    max(float(current_box[3]), float(word.get("y1"))),
                ]

        if current_box:
            rect = _pdf_scaled_rect(*current_box, page_width, page_height, page_number)
            if rect:
                line_rects.append(rect)

        return line_rects[:PDF_HIGHLIGHT_MAX_RECTS]

    def _load_pdf_chunks_with_geometry(self, file_path: str, file_name: str) -> list[Document]:
        """Load a PDF into chunks that retain page-level text highlight geometry."""
        if fitz is None:
            logger.warning("PyMuPDF is unavailable; PDF source highlights will fall back to page-level focus only.")
            return []

        documents: list[Document] = []
        with fitz.open(file_path) as pdf_document:
            for page_index, page in enumerate(pdf_document):
                page_number = page_index + 1
                page_text, words, page_width, page_height = self._extract_pdf_page_text(page, page_number)
                if not page_text:
                    continue

                search_cursor = 0
                for chunk_text in self.splitter.split_text(page_text):
                    chunk = str(chunk_text or "").strip()
                    if not chunk:
                        continue

                    start = page_text.find(chunk, search_cursor)
                    if start < 0:
                        start = page_text.find(chunk)
                    if start < 0:
                        start = search_cursor
                    end = min(len(page_text), start + len(chunk))
                    chunk_overlap = int(getattr(self.splitter, "_chunk_overlap", 50) or 0)
                    search_cursor = max(start + len(chunk) - chunk_overlap, start + 1)

                    rects = self._pdf_rects_for_chunk(words, start, end, page_width, page_height, page_number)
                    bounding_rect = _union_pdf_rects(rects)
                    metadata: dict[str, Any] = {"file_name": file_name, "page": page_index}

                    if rects and bounding_rect:
                        highlight_position = {
                            "boundingRect": bounding_rect,
                            "rects": rects,
                        }
                        metadata.update({
                            "highlight_position_json": json.dumps(highlight_position, ensure_ascii=False, sort_keys=True),
                            "text_quote": _compact_text(chunk, PDF_HIGHLIGHT_MAX_SOURCE_CHARS),
                        })

                    documents.append(Document(page_content=chunk, metadata=metadata))

        return documents
    
    def _load_file(self, file_path: str, file_name: str) -> list:
        """Load a file into LangChain documents based on file extention

        Args:
            file_path (str): path to file
            file_name (str): name of file

        Returns:
            list: returns the file content loader
        """
        ext = os.path.splitext(file_name)[-1].lower()
        if ext == ".pdf":
            loader = PyPDFLoader(file_path)
        elif ext == ".txt":
            loader = TextLoader(file_path)
        elif ext == ".docx":
            loader = Docx2txtLoader(file_path)
        elif ext == ".pptx":
            return self._load_pptx_file(file_path, file_name)
        else:
            raise ValueError(f"Unsupported file type: {ext}")
        docs = loader.load()
        return docs
    
    def load_file_content(self, file_path: str, file_name: str) -> str:
        """Load a file and return its raw text content
        Used when file is uploaded directly (no persistence to db)

        Args:
            file_path (str): file path
            file_name (str): file name

        Returns:
            str: raw file text content
        """
        docs = self._load_file(file_path, file_name)
        return "\n\n".join(doc.page_content for doc in docs)
    
    def load_file_content_from_bytes(self, file_bytes: bytes, file_name: str):
        """
        Write file bytes into a temp file, and using the tmp path to call load_file_content
        
        Args:
            file_bytes (bytes): raw file bytes
            file_name (str): file name used to get file extention to use correct file reader
        
        Returns:
            str: result of load_file_content
        """
        # print("---Processing Uploaded File---")
        logger.debug(f"Processing uploaded file: {file_name}")
        suffix = os.path.splitext(file_name or "")[-1]
        tmp = tempfile.NamedTemporaryFile(delete=False, suffix=suffix)
        try:
            tmp.write(file_bytes)
            tmp.close()
            return self.load_file_content(tmp.name, file_name)
        finally:
            tmp.close()
            os.remove(tmp.name)
    
    def _file_source_metadata(self, room_id: str, item: Any, display_name: str, url: str, doc: Document) -> dict[str, Any]:
        """Build flat metadata and a typed source ref for one file-derived page/slide."""
        source_ref = dict(_payload_get(item, "source_ref", {}) or {})
        metadata = dict(_payload_get(item, "metadata", {}) or {})
        source_type = str(_payload_get(item, "source_type", "resource") or "resource")
        resource_id = str(source_ref.get("resourceId") or metadata.get("resource_id") or _payload_get(item, "id", "") or "")
        page = doc.metadata.get("page")
        slide = doc.metadata.get("slide")

        if page is not None and "pageNumber" not in source_ref:
            try:
                source_ref["pageNumber"] = int(page) + 1
            except (TypeError, ValueError):
                pass
        if slide is not None and "slideNumber" not in source_ref:
            try:
                source_ref["slideNumber"] = int(slide)
            except (TypeError, ValueError):
                pass

        source_ref = {
            "type": source_ref.get("type") or source_type,
            "roomId": source_ref.get("roomId") or room_id,
            "label": _source_label(source_ref, display_name),
            "resourceId": resource_id,
            "sourceId": source_ref.get("sourceId") or resource_id or display_name,
            **source_ref,
        }

        return _clean_metadata({
            **metadata,
            "file_name": display_name,
            "label": source_ref["label"],
            "room_id": room_id,
            "source_type": source_ref["type"],
            "source_id": source_ref.get("sourceId") or resource_id or display_name,
            "resource_id": resource_id,
            "source_url": url,
            "page": source_ref.get("pageNumber") or "",
            "slide": source_ref.get("slideNumber") or "",
            "highlight_position_json": doc.metadata.get("highlight_position_json") or "",
            "text_quote": doc.metadata.get("text_quote") or "",
            "source_ref_json": json.dumps(source_ref, ensure_ascii=False, sort_keys=True),
        })

    def _domain_document_metadata(self, room_id: str, item: Any) -> dict[str, Any]:
        """Build flat metadata and source ref storage for one app-owned Domain record."""
        source_ref = dict(_payload_get(item, "source_ref", {}) or {})
        metadata = dict(_payload_get(item, "metadata", {}) or {})
        source_type = str(_payload_get(item, "source_type", source_ref.get("type") or "domain") or "domain")
        source_id = str(_payload_get(item, "id", "") or source_ref.get("sourceId") or "")
        title = str(_payload_get(item, "title", "") or source_ref.get("label") or source_id or "Domain source")

        source_ref = {
            "type": source_ref.get("type") or source_type,
            "roomId": source_ref.get("roomId") or room_id,
            "label": _source_label(source_ref, title),
            "sourceId": source_ref.get("sourceId") or source_id,
            **source_ref,
        }

        return _clean_metadata({
            **metadata,
            "file_name": source_ref["label"],
            "label": source_ref["label"],
            "title": title,
            "room_id": room_id,
            "source_type": source_ref["type"],
            "source_id": source_ref.get("sourceId") or source_id,
            "resource_id": source_ref.get("resourceId") or metadata.get("resource_id") or "",
            "message_id": source_ref.get("messageId") or metadata.get("message_id") or "",
            "annotation_id": source_ref.get("annotationId") or metadata.get("annotation_id") or "",
            "session_id": source_ref.get("sessionId") or metadata.get("session_id") or "",
            "poll_id": source_ref.get("pollId") or metadata.get("poll_id") or "",
            "channel": source_ref.get("channel") or metadata.get("channel") or "",
            "date": source_ref.get("date") or metadata.get("date") or "",
            "starts_at": source_ref.get("startsAt") or metadata.get("starts_at") or "",
            "ends_at": source_ref.get("endsAt") or metadata.get("ends_at") or "",
            "start_ts": metadata.get("start_ts") or _iso_to_epoch_ms(source_ref.get("startsAt") or metadata.get("starts_at")),
            "end_ts": metadata.get("end_ts") or _iso_to_epoch_ms(source_ref.get("endsAt") or metadata.get("ends_at")),
            "range_start_ts": metadata.get("range_start_ts") or _iso_to_epoch_ms(source_ref.get("startsAt") or metadata.get("starts_at")),
            "range_end_ts": metadata.get("range_end_ts") or _iso_to_epoch_ms(source_ref.get("endsAt") or metadata.get("ends_at")),
            "kind": source_ref.get("kind") or metadata.get("kind") or "",
            "source_ref_json": json.dumps(source_ref, ensure_ascii=False, sort_keys=True),
        })

    async def _embed_file_documents(self, room_id: str, urls: list) -> dict[str, Any]:
        """Load, chunk, and store file documents without clearing existing room records."""
        results = {
            "success": [],
            "failed": [],
            "total_chunks": 0,
        }
        
        async with httpx.AsyncClient(base_url=NODE_BASE_URL, timeout=30) as client:
            for item in urls:
                url = str(_payload_get(item, "url", item if isinstance(item, str) else "") or "")
                display_name = str(_payload_get(item, "file_name", "") or url.split("/")[-1])

                url_file_name = url.split("/")[-1]
                suffix = os.path.splitext(display_name or url_file_name)[-1]
                
                try:
                    logger.debug(f"Fetching file: {display_name} from {url}")
                    response = await client.get(url)
                    response.raise_for_status()
                    
                    # Save to temp file for loader
                    with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as tmp:
                        tmp.write(response.content)
                        tmp_path = tmp.name
                    
                    try:
                        is_pdf = suffix.lower() == ".pdf" or display_name.lower().endswith(".pdf")
                        chunks = []

                        if is_pdf:
                            chunks = self._load_pdf_chunks_with_geometry(tmp_path, display_name)

                        if chunks:
                            for chunk in chunks:
                                chunk.metadata = self._file_source_metadata(room_id, item, display_name, url, chunk)
                        else:
                            docs = self._load_file(tmp_path, display_name)
                            for doc in docs:
                                doc.metadata = self._file_source_metadata(room_id, item, display_name, url, doc)
                            chunks = self.splitter.split_documents(docs)

                        if chunks:
                            self.db.add_documents(chunks)
                        
                        results["success"].append(display_name)
                        results["total_chunks"] += len(chunks)
                        logger.debug(f"Embedded {len(chunks)} chunks from: {display_name}")
                    finally:
                        os.remove(tmp_path)
                except Exception as e:
                    logger.error(f"Failed to embed file {display_name}: {e}", exc_info=True)
                    results["failed"].append({"file": display_name, "error": str(e)})

        return results

    async def embed_room_documents(self, room_id: str, urls: list) -> int:
        """Load, chunk and store document files in ChromaDB for legacy /embed callers."""
        logger.info(f"Embedding {len(urls)} document(s) for room: {room_id}")
        self.clear(room_id = room_id)
        results = await self._embed_file_documents(room_id, urls)
        logger.info(f"Embed complete for room {room_id} | chunks: {results['total_chunks']} | failed: {len(results['failed'])}")
        return results

    async def sync_domain_corpus(self, room_id: str, files: list, documents: list) -> dict[str, Any]:
        """Replace a room corpus with files plus typed Domain records from the app server."""
        logger.info(
            f"Domain corpus sync | room: {room_id} | files: {len(files)} | records: {len(documents)}"
        )
        self.clear(room_id=room_id)
        results = await self._embed_file_documents(room_id, files)

        for item in documents:
            title = str(_payload_get(item, "title", "") or _payload_get(item, "id", "Domain source"))
            text = str(_payload_get(item, "text", "") or "").strip()
            if not text:
                results["failed"].append({"file": title, "error": "Domain record did not contain text"})
                continue

            try:
                doc = Document(
                    page_content=text,
                    metadata=self._domain_document_metadata(room_id, item),
                )
                chunks = self.splitter.split_documents([doc])
                if chunks:
                    self.db.add_documents(chunks)
                results["success"].append(title)
                results["total_chunks"] += len(chunks)
            except Exception as error:
                logger.error(f"Failed to embed Domain record {title}: {error}", exc_info=True)
                results["failed"].append({"file": title, "error": str(error)})

        logger.info(
            f"Domain corpus sync complete | room: {room_id} | chunks: {results['total_chunks']} | failed: {len(results['failed'])}"
        )
        return results
    
    def _search_filter(self, room_id: str, scope: Optional[dict[str, Any]] = None) -> dict[str, Any]:
        """Translate optional tool scope into a Chroma metadata filter."""
        scope = scope or {}
        clauses: list[dict[str, Any]] = [{"room_id": room_id}]
        raw_source_type = str(scope.get("source_type") or "").strip().lower()
        source_type = SOURCE_TYPE_ALIASES.get(raw_source_type, raw_source_type)

        if isinstance(source_type, list):
            clauses.append({"source_type": {"$in": source_type}})
        elif source_type:
            clauses.append({"source_type": source_type})

        for scope_key, metadata_key in (
            ("source_id", "source_id"),
            ("resource_id", "resource_id"),
            ("message_id", "message_id"),
            ("annotation_id", "annotation_id"),
            ("session_id", "session_id"),
            ("poll_id", "poll_id"),
            ("channel", "channel"),
            ("date", "date"),
        ):
            value = str(scope.get(scope_key) or "").strip()
            if value:
                clauses.append({metadata_key: value})

        date_from_ts = int(scope.get("date_from_ts") or _iso_to_epoch_ms(scope.get("date_from")) or 0)
        date_to_ts = int(scope.get("date_to_ts") or _iso_to_epoch_ms(scope.get("date_to")) or 0)
        if date_from_ts or date_to_ts:
            temporal_fields = ["start_ts", "range_start_ts"]
            temporal_options = []
            for field in temporal_fields:
                field_clauses = []
                if date_from_ts:
                    field_clauses.append({field: {"$gte": date_from_ts}})
                if date_to_ts:
                    field_clauses.append({field: {"$lte": date_to_ts}})
                if len(field_clauses) == 1:
                    temporal_options.append(field_clauses[0])
                elif field_clauses:
                    temporal_options.append({"$and": field_clauses})

            if len(temporal_options) == 1:
                clauses.append(temporal_options[0])
            elif temporal_options:
                clauses.append({"$or": temporal_options})

        return clauses[0] if len(clauses) == 1 else {"$and": clauses}

    def _scope_source_types(self, scope: Optional[dict[str, Any]] = None) -> list[str]:
        """Resolve a tool scope into concrete source_type metadata values."""
        raw_source_type = str((scope or {}).get("source_type") or "").strip().lower()
        source_type = SOURCE_TYPE_ALIASES.get(raw_source_type, raw_source_type)
        if isinstance(source_type, list):
            return source_type
        return [source_type] if source_type else []

    def _should_use_metadata_first(self, scope: Optional[dict[str, Any]] = None) -> bool:
        """Use exact metadata lookup for structured Coordidate list/range queries.

        Calendar questions such as "any meetings coming up?" are not primarily
        semantic search problems. They should first retrieve records matching
        typed source and time metadata, then let the model summarize them.
        """
        scope = scope or {}
        source_types = self._scope_source_types(scope)
        has_coordidate_scope = any(
            source_type in {"coordidate_session", "coordidate_poll"} for source_type in source_types
        )
        has_temporal_scope = any(
            scope.get(key)
            for key in ("timeframe", "date", "date_from", "date_to", "date_from_ts", "date_to_ts")
        )
        return has_coordidate_scope and has_temporal_scope

    def _metadata_filter_documents(self, metadata_filter: dict[str, Any], k: int, scope: Optional[dict[str, Any]] = None) -> Optional[list[Document]]:
        """Fetch metadata-filtered records directly from Chroma without similarity scoring."""
        try:
            raw = self.db._collection.get(
                where=metadata_filter,
                limit=max(k, 20),
                include=["documents", "metadatas"],
            )
        except Exception as error:
            logger.warning(f"Metadata-first Chroma lookup failed, falling back to vector search: {error}")
            return None

        documents = []
        for page_content, metadata in zip(raw.get("documents") or [], raw.get("metadatas") or []):
            if not page_content:
                continue
            document = Document(page_content=page_content, metadata=dict(metadata or {}))
            document.metadata["relevance_score"] = 1.0
            documents.append(document)

        reverse = str((scope or {}).get("timeframe") or "").lower() in {"past", "previous"}

        def temporal_sort_key(document: Document) -> int:
            metadata = document.metadata or {}
            return int(metadata.get("start_ts") or metadata.get("range_start_ts") or 0)

        documents.sort(key=temporal_sort_key, reverse=reverse)
        return documents[:k]

    def source_ref_for_document(self, document: Document) -> dict[str, Any]:
        """Rehydrate the typed source ref stored with a retrieved vector chunk."""
        metadata = document.metadata or {}
        raw_ref = metadata.get("source_ref_json") or ""
        try:
            source_ref = json.loads(raw_ref) if raw_ref else {}
        except json.JSONDecodeError:
            source_ref = {}

        source_ref.setdefault("type", metadata.get("source_type") or "resource")
        source_ref.setdefault("roomId", metadata.get("room_id") or "")
        source_ref.setdefault("sourceId", metadata.get("source_id") or "")
        source_ref.setdefault("resourceId", metadata.get("resource_id") or "")
        source_ref.setdefault("messageId", metadata.get("message_id") or "")
        source_ref.setdefault("annotationId", metadata.get("annotation_id") or "")
        source_ref.setdefault("sessionId", metadata.get("session_id") or "")
        source_ref.setdefault("pollId", metadata.get("poll_id") or "")
        source_ref.setdefault("channel", metadata.get("channel") or "")
        source_ref.setdefault("label", metadata.get("label") or metadata.get("file_name") or "Domain source")
        source_ref.setdefault("startsAt", metadata.get("starts_at") or "")
        source_ref.setdefault("endsAt", metadata.get("ends_at") or "")
        if metadata.get("page") and not source_ref.get("pageNumber"):
            source_ref["pageNumber"] = metadata.get("page")
        if metadata.get("slide") and not source_ref.get("slideNumber"):
            source_ref["slideNumber"] = metadata.get("slide")
        raw_highlight_position = metadata.get("highlight_position_json") or ""
        if raw_highlight_position and not source_ref.get("highlightPosition"):
            try:
                highlight_position = _valid_highlight_position(json.loads(raw_highlight_position))
                if highlight_position:
                    source_ref["highlightPosition"] = highlight_position
            except (TypeError, json.JSONDecodeError):
                logger.debug("Ignoring invalid PDF highlight metadata on retrieved source", exc_info=True)
        if metadata.get("text_quote") and not source_ref.get("textQuote"):
            source_ref["textQuote"] = str(metadata.get("text_quote") or "")[:PDF_HIGHLIGHT_MAX_SOURCE_CHARS]
        if metadata.get("relevance_score") is not None:
            source_ref["score"] = metadata.get("relevance_score")
        source_ref["snippet"] = _compact_text(document.page_content)
        return {key: value for key, value in source_ref.items() if value not in (None, "")}

    def search(self, query: str, room_id: str, k: int = 5, scope: Optional[dict[str, Any]] = None) -> list:
        """Search ChromaDB for relevant chunks.

        Args:
            query (str): the query to search documents based off
            room_id (str): room id of documents to retrieve.
            k (int, optional): Number of chunks to retrieve. Defaults to 5.
            scope (dict, optional): metadata filters such as source_type, channel, or resource_id.

        Returns:
            list: list of chunks retrieved
        """
        # print("---Search for Chunks---")
        # print("Chunk count:", self.db._collection.count())
        metadata_filter = self._search_filter(room_id, scope)
        logger.debug(
            f"Search | room: {room_id} | query: {query!r} | filter: {metadata_filter} | total chunks in db: {self.db._collection.count()}"
        )

        if self._should_use_metadata_first(scope):
            metadata_results = self._metadata_filter_documents(metadata_filter, k, scope)
            if metadata_results is not None:
                logger.debug(f"Metadata-first search results | found: {len(metadata_results)}")
                return metadata_results
    
        try:
            scored_results = self.db.similarity_search_with_relevance_scores(
                query=query,
                k=k,
                filter=metadata_filter,
            )
            
            best_score = max((float(score) for _, score in scored_results), default=0.0)
            # Keep documents that are both absolutely relevant and close enough to
            # the best match. This avoids citing weak side hits that share a keyword
            # but do not actually answer the question.
            relevance_floor = max(SEARCH_MIN_RELEVANCE, best_score * 0.7)

            filtered_results = []
            for document, score in scored_results:
                document.metadata["relevance_score"] = float(score)
                if score >= relevance_floor:
                    filtered_results.append(document)

            logger.debug(f"Search results | found: {len(scored_results)} | after filter: {len(filtered_results)} | best core: {best_score:.3f} | floor: {relevance_floor:.3f}")
            return filtered_results
        
        except Exception as error:
            logger.warning(f"Revelance search failed, falling back to similarity search: {error}")
            documents = self.db.similarity_search(
                query=query,
                k=k,
                filter=metadata_filter,
            )
            return documents
    
    def clear(self, room_id: str):
        """Clear documents from ChromaDB
        
        If room_id is provided only clear that rooms documents

        Args:
            room_id (str): filter of room id to remove documents from.
        """
        logger.info(f"Clearning corpus for room: {room_id}")
        self.db.delete(where={"room_id": room_id})
